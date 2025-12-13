import streamlit as st
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
import io
import json
import re
import random
import zipfile
import time
import os

# --- 页面配置 ---
st.set_page_config(
    page_title="京东直播手卡生成器 Web版",
    page_icon="⚡",
    layout="wide"
)

# --- 工具函数 ---

# 伪装浏览器头
def get_headers():
    user_agents = [
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36",
        "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
    ]
    return {
        "User-Agent": random.choice(user_agents),
        "Referer": "https://item.jd.com/",
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8",
        "Accept-Language": "zh-CN,zh;q=0.9"
    }

def scrape_jd_sku(sku):
    """抓取京东商品标题和主图"""
    url = f"https://item.jd.com/{sku}.html"
    info = {"sku": sku, "title": "", "image_url": ""}
    
    try:
        r = requests.get(url, headers=get_headers(), timeout=10)
        r.encoding = r.apparent_encoding
        soup = BeautifulSoup(r.text, 'html.parser')
        
        # 1. 抓标题
        raw_title = ""
        title_tag = soup.select_one("div.sku-name")
        if title_tag: raw_title = title_tag.get_text(strip=True)
        if not raw_title and soup.title: raw_title = soup.title.string.split('-')[0].strip()
        
        if raw_title:
            # 清理标题中的换行和多余空格，防止干扰 Prompt
            clean_title = raw_title.replace("京东", "").replace("自营", "").replace("\n", " ").strip()
            info["title"] = clean_title
        else:
            info["title"] = f"商品_{sku}" # 如果这里由于反爬没抓到标题，AI就会生成垃圾内容

        # 2. 抓主图
        candidates = []
        img_tag = soup.select_one("#spec-img")
        if img_tag:
            candidates.append(img_tag.get('data-origin'))
            candidates.append(img_tag.get('src'))
        
        # 正则补充匹配
        patterns = re.findall(r'//img\d{1,2}\.360buyimg\.com/n[01]/jfs/[^"]+\.jpg', r.text)
        candidates.extend(patterns)

        for img in candidates:
            if img and "jfs" in img and ".jpg" in img:
                if not img.startswith("http"):
                    img = "https:" + img if img.startswith("//") else "https://" + img
                # 替换为高清大图
                img = img.replace("/n1/", "/n0/").replace("/n5/", "/n0/")
                info["image_url"] = img
                break
                
        return info
    except Exception as e:
        # 不直接报错，返回基础信息，避免打断循环
        print(f"SKU {sku} 抓取异常: {e}")
        info["title"] = f"商品_{sku}" 
        return info

def download_image_to_memory(url):
    """下载图片到内存字节流"""
    if not url: return None
    try:
        r = requests.get(url, headers=get_headers(), timeout=10, verify=False)
        return io.BytesIO(r.content)
    except Exception as e:
        st.error(f"图片下载失败: {e}")
        return None

def call_ai_generate_points(product_name, api_key, base_url):
    """调用 AI 生成卖点 (User/System 分离版 + 强制字数约束)"""
    if not api_key:
        return {"selling_point_1": "请填写API Key", "selling_point_2": "以生成智能卖点"}

    # --- 安全检查 ---
    if product_name.startswith("商品_") and product_name[3:].isdigit():
        return {
            "selling_point_1": "标题抓取失败",
            "selling_point_2": "无法生成智能卖点",
            "selling_point_3": "请检查 SKU 是否正确",
            "selling_point_4": "或稍微放慢生成速度"
        }

    headers = {"Content-Type": "application/json", "Authorization": f"Bearer {api_key}"}
    
    # --- 核心修改：将提示词拆分为 System（规则）和 User（内容） ---
    
    # System Prompt: 永远不变的法则，强制 AI 每次都遵守
    system_prompt = """
    你是一名带货过亿的金牌直播运营。你的唯一任务是为商品撰写极具煽动性的【直播手卡文案】。
    
    【必须遵守的死命令】：
    1. **拒绝简短**：每一条卖点的【详细解释】部分必须在 30 个字以上，严禁只写一句话。
    2. **拒绝通用**：禁止使用“性价比高”、“非常好用”等万金油词汇，必须结合商品名称中的具体参数（如材质、功率、成分）来发挥。
    3. **结构要求**：先写痛点（没有它多麻烦），再写爽点（有了它多舒服）。
    4. **语气要求**：口语化、紧迫感、像是朋友按头安利。
    
    【输出格式】：
    必须返回标准 JSON 对象，Key 必须严格为 selling_point_1, selling_point_2, selling_point_3, selling_point_4。
    """

    # User Prompt: 每次变化的商品
    user_prompt = f"""
    需生成的商品名称：【{product_name}】
    
    请立刻执行上述规则，生成 4 个详细的高转化卖点。
    """
    
    data = {
        "model": "deepseek-chat", 
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        "temperature": 0.85, # 稍微提高一点随机性，避免内容重复
        "response_format": {"type": "json_object"},
        # 添加随机种子，防止 API 缓存导致的重复或偷懒
        "seed": random.randint(1, 10000) 
    }

    try:
        resp = requests.post(f"{base_url}/chat/completions", headers=headers, json=data, timeout=30)
        result = resp.json()
        if 'error' in result:
            st.error(f"AI 接口报错: {result['error']['message']}")
            return {}
        content = result['choices'][0]['message']['content']
        return json.loads(content)
    except Exception as e:
        st.error(f"AI 请求异常: {e}")
        return {}

def process_ppt(template_file_obj, data_list):
    """批量生成 PPT 并打包成 ZIP"""
    zip_buffer = io.BytesIO()
    
    with zipfile.ZipFile(zip_buffer, "a", zipfile.ZIP_DEFLATED, False) as zip_file:
        for data in data_list:
            # 每次都需要重新加载模板（指针归零）
            template_file_obj.seek(0)
            prs = Presentation(template_file_obj)
            slide = prs.slides[0]

            # 文本替换函数
            def replace_text(name, text):
                for shape in slide.shapes:
                    if shape.name == name and shape.has_text_frame:
                        shape.text_frame.text = str(text)
                        return
                    if shape.shape_type == 6: # Group
                        for sub in shape.shapes:
                            if sub.name == name and sub.has_text_frame:
                                sub.text_frame.text = str(text)
                                return

            # 执行替换
            replace_text("product_name", data['title'])
            replace_text("product_sku", data['sku'])
            replace_text("price_live", data['price'])
            
            points = data.get('points', {})
            replace_text("selling_point_1", points.get('selling_point_1', ''))
            replace_text("selling_point_2", points.get('selling_point_2', ''))
            replace_text("selling_point_3", points.get('selling_point_3', ''))
            replace_text("selling_point_4", points.get('selling_point_4', ''))

            # 图片替换
            if data['image_bytes']:
                found_img = False
                for shape in slide.shapes:
                    if shape.name == "product_image":
                        left, top, width, height = shape.left, shape.top, shape.width, shape.height
                        # 移除旧图
                        sp = shape._element
                        sp.getparent().remove(sp)
                        # 添加新图
                        slide.shapes.add_picture(data['image_bytes'], left, top, width, height)
                        found_img = True
                        break
            
            # 保存单个 PPT 到内存
            ppt_buffer = io.BytesIO()
            prs.save(ppt_buffer)
            # 添加到 ZIP
            zip_file.writestr(f"{data['sku']}.pptx", ppt_buffer.getvalue())
    
    return zip_buffer

# --- UI 布局 ---

st.title("⚡ 京东直播手卡全自动生成器 (Web版)")
st.markdown("上传 PPT 模板，输入 SKU，自动抓取信息 + AI 生成痛点卖点，一键导出 PPT。")

with st.sidebar:
    st.header("🧠 1. AI 配置")
    api_key = st.text_input("API Key", type="password", help="推荐使用 DeepSeek API")
    base_url = st.text_input("Base URL", value="https://api.deepseek.com")
    st.info("如果没有 Key，卖点部分将为空，但基础信息仍会生成。")
    
    st.divider()
    st.header("📂 2. 模板设置")
    
    # --- 模板加载逻辑 ---
    uploaded_template = st.file_uploader("上传 .pptx 模板文件 (可选)", type=["pptx"])
    
    # 默认模板文件名
    DEFAULT_TEMPLATE_NAME = "template.pptx"
    
    final_template_file = None
    
    if uploaded_template:
        st.success(f"✅ 使用上传的模板: {uploaded_template.name}")
        final_template_file = uploaded_template
    elif os.path.exists(DEFAULT_TEMPLATE_NAME):
        st.info(f"ℹ️ 未上传模板，将使用系统默认模板 ({DEFAULT_TEMPLATE_NAME})")
        # 将本地文件读入内存，模拟 uploaded_file 的行为
        with open(DEFAULT_TEMPLATE_NAME, "rb") as f:
            final_template_file = io.BytesIO(f.read())
    else:
        st.warning(f"⚠️ 请上传模板！(且未在服务器找到默认模板 {DEFAULT_TEMPLATE_NAME})")

    st.markdown("---")
    st.caption("**模板制作说明**：PPT中需包含以下元素名称（Selection Pane）：\n`product_name`, `product_sku`, `price_live`, `product_image`, `selling_point_1`~`4`")

st.header("📝 3. 商品与价格")
col1, col2 = st.columns([3, 1])

with col1:
    # --- 批量输入逻辑说明 ---
    st.markdown("**输入 SKU 和 价格** (格式：`SKU, 价格`，一行一个)")
    sku_input = st.text_area(
        "SKU列表", 
        height=180, 
        placeholder="例如：\n1000123456, 9.9\n1000888888, 19.9\n1000999999 (未填价格将使用右侧默认价)"
    )

with col2:
    default_price = st.text_input("默认兜底价格", value="待定")
    st.caption("如果左侧某一行只写了 SKU 没写价格，将自动使用此价格。")

# --- 执行逻辑 ---

if st.button("🚀 开始生成", type="primary", use_container_width=True):
    if not final_template_file:
        st.error("❌ 无法开始：没有可用的 PPT 模板（请上传或联系管理员添加默认模板）。")
    elif not sku_input.strip():
        st.error("❌ 请输入至少一个 SKU！")
    else:
        # 1. 解析 SKU 和 价格
        lines = sku_input.strip().split('\n')
        tasks = []
        
        for line in lines:
            line = line.strip()
            if not line: continue
            
            # 兼容中文逗号
            line = line.replace('，', ',')
            
            parts = line.split(',')
            current_sku = parts[0].strip()
            
            # 如果有逗号分隔，取第二个作为价格；否则使用默认价格
            current_price = parts[1].strip() if len(parts) > 1 else default_price
            
            if current_sku:
                tasks.append({"sku": current_sku, "price": current_price})

        if not tasks:
            st.error("❌ 未识别到有效 SKU。")
            st.stop()
            
        processed_data = []
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        # 2. 循环处理
        for idx, task in enumerate(tasks):
            sku = task['sku']
            price = task['price']
            
            status_text.text(f"正在处理 ({idx+1}/{len(tasks)}): SKU {sku} ...")
            
            # --- 增加随机延时，防止反爬 ---
            if idx > 0:
                sleep_time = random.uniform(1.5, 4.0)
                time.sleep(sleep_time)

            # 抓取
            info = scrape_jd_sku(sku)
            if not info:
                # 即使抓取失败也可以跳过，或者生成一个空的占位
                continue
                
            info['price'] = price
            
            # AI 生成
            if api_key:
                info['points'] = call_ai_generate_points(info['title'], api_key, base_url)
            else:
                info['points'] = {}
            
            # 下载图片
            info['image_bytes'] = download_image_to_memory(info['image_url'])
            
            processed_data.append(info)
            progress_bar.progress((idx + 1) / len(tasks))
            
        status_text.text("正在生成 PPT 文件...")
        
        # 3. 生成 PPT 压缩包
        if processed_data:
            try:
                zip_io = process_ppt(final_template_file, processed_data)
                
                st.success(f"🎉 成功生成 {len(processed_data)} 个手卡！")
                
                st.download_button(
                    label="📥 下载所有手卡 (ZIP压缩包)",
                    data=zip_io.getvalue(),
                    file_name="Live_Cards_Output.zip",
                    mime="application/zip",
                    type="primary"
                )
            except Exception as e:
                st.error(f"生成 PPT 时发生错误 (可能是模板格式问题): {e}")
        else:
            st.error("未能生成有效数据，请检查 SKU 是否正确。")